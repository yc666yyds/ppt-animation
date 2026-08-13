# -*- coding: utf-8 -*-
import os, re, zipfile, xml.etree.ElementTree as ET, tempfile, shutil
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.enum.shapes import MSO_SHAPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
try:
    import win32com.client
    HAS_WIN32COM = True
except ImportError:
    HAS_WIN32COM = False
OUTPUT_DIR = r'D:\Workspace\PPT成品'
os.makedirs(OUTPUT_DIR, exist_ok=True)
SW = Inches(13.333); SH = Inches(7.5)
BG = RGBColor(0x0D, 0x1B, 0x2A)
CD = RGBColor(0x1B, 0x3A, 0x5C)
GD = RGBColor(0xE8, 0xB8, 0x30)
WH = RGBColor(0xFF, 0xFF, 0xFF)
LT = RGBColor(0xE0, 0xE0, 0xE0)
GY = RGBColor(0x88, 0x99, 0xAA)
# COM 转场效果 ID 映射（已通过测试验证）
COM_TRANSITION_IDS = {
    'fade': 3849,
    'wipe': 2817,
    'split': 3585,
    'zoom': 3073,
    'push': 1285,
    'cover': 3853,
    'dissolve': 1281,
    'blind': 1284,
    'checkerboard': 3852,
    'coin': 1282,
    'ripple': 2818,
    'starburst': 1281,
    'sunbeam': 2817,
    'triangle': 3073,
    'wedge': 1285,
    'wheel': 3853,
    'zipper': 3852,
    'diamond': 3846,
    'random': 3848,
}
TRANS_MAP = {
    'fade': 2, 'wipe': 7, 'split': 13, 'zoom': 41, 
    'push': 10, 'cover': 6,
    # 扩展转场
    'dissolve': 41, 'blind': 8, 'checkerboard': 9,
    'coin': 11, 'diamond': 15, 'random': 16,
    'ripple': 18, 'starburst': 20, 'sunbeam': 22,
    'triangle': 23, 'wedge': 24, 'wheel': 25, 'zipper': 26
}
TRANSITION_XML = {
    'fade': '<p:fade/>', 'wipe': '<p:wipe/>', 'split': '<p:split/>',
    'zoom': '<p:zoom/>', 'push': '<p:push/>', 'cover': '<p:cover/>',
    # 扩展转场
    'dissolve': '<p:dissolve/>', 'blind': '<p:blind/>',
    'checkerboard': '<p:checkerboard/>', 'coin': '<p:coin/>',
    'diamond': '<p:diamond/>', 'random': '<p:random/>',
    'ripple': '<p:ripple/>', 'starburst': '<p:starburst/>',
    'sunbeam': '<p:sunbeam/>', 'triangle': '<p:triangle/>',
    'wedge': '<p:wedge/>', 'wheel': '<p:wheel/>', 'zipper': '<p:zipper/>'
}
EFFECT_MAP = {'appear':1,'checkerboard':2,'circle':86,'box':86,'spin':3,'fly':4,'blend':17,'blur':28,'compress':40,'dissolve':41,'explode':42,'fade':43,'glow':35,'grow':53,'misty':59,'ripple':64,'reveal':65,'roll':66,'shrink':71,'swizzle':100,'teeter':105,'typeWriter':106}

# COM entrance ID 映射（XML presetID -> COM msoAnimEffect ID）
# 通过扫描PowerPoint COM枚举验证，COM ID和XML presetID不是一一对应关系
# COM entrance ID映射（XML presetID -> COM msoAnimEffect ID）
# 通过PowerPoint COM枚举扫描验证
# 注意：部分效果在COM中不存在entrance变体，需要降级到XML模式
COM_ENTRANCE_IDS = {
    1: 1,     # appear
    2: 2,     # checkerboard
    3: 3,     # spin
    4: 4,     # fly
    7: 7,     # wipe
    17: 17,   # blend
    28: 28,   # blur
    35: 35,   # glow
    37: 40,   # compress
    39: 42,   # explode
    41: 41,   # dissolve
    43: 43,   # fade
    51: 53,   # grow
    # 以下效果在COM中没有正确的entrance实现，标记为None表示需要降级XML
    59: None, # misty - COM中只有emph变体
    64: None, # ripple - COM中只有emph变体
    65: None, # reveal - COM中只有emph变体
    66: None, # roll - COM中只有emph变体
    71: None, # shrink - COM中只有emph变体
    86: None, # circle/box - COM中只有path变体
    100: None,# swizzle - COM中只有path变体
    105: None,# teeter - COM中只有path变体
    106: None,# typeWriter - COM中只有path变体
}
# 反向映射：XML presetID -> COM ID（排除None）
_XML_TO_COM_ID = {v: k for k, v in COM_ENTRANCE_IDS.items() if v is not None}
_XML_TO_COM_ID = {v: k for k, v in COM_ENTRANCE_IDS.items()}
EFFECT_APPEAR=1; EFFECT_CHECKERBOARD=2; EFFECT_CIRCLE=86; EFFECT_BOX=86

# ============ 默认动画模板 ============
# 简洁商务风格默认模板
DEFAULT_TRANSITIONS = ['fade', 'zoom', 'fade']  # 转场：淡入淡出 + 缩放（简洁有层次）
DEFAULT_EFFECTS = ['appear', 'fly', 'fade', 'grow', 'wipe', 'blend']  # 入场：简洁多样

# 动画风格模板（名称 -> effects列表）
EFFECT_STYLES = {
    '简洁商务': ['appear', 'fly', 'fade', 'grow', 'wipe', 'blend'],
    '动态活力': ['fly', 'spin', 'explode', 'zoom', 'roll', 'teeter'],
    '优雅柔和': ['fade', 'blur', 'dissolve', 'grow', 'reveal', 'misty'],
    '科技未来': ['spin', 'glow', 'ripple', 'swizzle', 'blur', 'zoom'],
}

# 转场风格模板
TRANSITION_STYLES = {
    '简洁商务': ['fade', 'zoom', 'fade'],
    '动态活力': ['zoom', 'wipe', 'split', 'starburst'],
    '优雅柔和': ['fade', 'dissolve', 'push', 'fade'],
    '科技未来': ['zoom', 'split', 'cover', 'ripple'],
}
from .colors import COLOR_SCHEMES, DEFAULT_SCHEME, get_color_scheme_names, get_color_scheme, apply_scheme

def get_animation_recommendations():
    """获取动画推荐列表"""
    recs = []
    for style, effects in EFFECT_STYLES.items():
        trans = TRANSITION_STYLES[style]
        recs.append(f"{style}: 转场={trans}, 入场={effects}")
    return recs


# 扩展动画效果映射（presetID -> 效果名称）
EFFECT_PRESETS = {
    1: 'appear',           # 出现
    2: 'checkerboard',     # 棋盘
    3: 'spin',             # 旋转
    4: 'fly',              # 飞入
    5: 'fade',             # 淡入
    6: 'glow',             # 发光
    7: 'wipe',             # 擦除
    17: 'blend',           # 混合
    28: 'blur',            # 模糊
    35: 'glow',            # 发光（变体）
    41: 'dissolve',        # 溶解
    42: 'explode',         # 爆炸
    43: 'fade',            # 淡入（变体）
    53: 'grow',            # 放大
    59: 'misty',           # 迷雾
    64: 'ripple',          # 波纹
    65: 'reveal',          # 揭示
    66: 'roll',            # 滚动
    71: 'shrink',          # 缩小
    86: 'circle',          # 圆形
    100: 'swizzle',        # 扭曲
    105: 'teeter',         # 摇晃
    106: 'typeWriter'      # 打字机
}

# 动画效果XML生成器
def generate_effect_xml(preset_id, spid, cid):
    """根据preset_id生成对应的动画XML"""
    subtype_map = {
        1: (0, ''),                      # appear
        2: (4, 'checkerboard_anim'),     # checkerboard
        3: (0, 'spin_anim'),             # spin
        4: (16, 'fly_anim'),             # fly
        5: (10, 'fade_anim'),            # fade
        6: (16, 'glow_anim'),            # glow
        7: (4, 'wipe_anim'),             # wipe
        17: (0, 'blend_anim'),           # blend
        28: (0, 'blur_anim'),            # blur
        35: (16, 'glow_anim'),           # glow (variant)
        41: (0, 'dissolve_anim'),        # dissolve
        42: (0, 'explode_anim'),         # explode
        43: (10, 'fade_anim'),           # fade (variant)
        53: (0, 'grow_anim'),            # grow
        59: (0, 'misty_anim'),           # misty
        64: (0, 'ripple_anim'),          # ripple
        65: (0, 'reveal_anim'),          # reveal
        66: (0, 'roll_anim'),            # roll
        71: (0, 'shrink_anim'),          # shrink
        86: (0, 'circle_anim'),          # circle
        100: (0, 'swizzle_anim'),        # swizzle
        105: (0, 'teeter_anim'),         # teeter
        106: (0, 'typeWriter_anim')      # typeWriter
    }
    
    if preset_id not in subtype_map:
        return 0, ''
    
    subtype, effect_type = subtype_map[preset_id]
    
    # 生成对应的动画效果XML
    if effect_type == '':
        extra = ''
    elif effect_type == 'checkerboard_anim':
        extra = f'<p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_x</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav></p:tavLst></p:anim><p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid+1}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="1+#ppt_h/2"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav></p:tavLst></p:anim>'
    elif effect_type == 'spin_anim':
        extra = f'<p:animEffect transition="in" filter="spin(in)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'fly_anim':
        extra = f'<p:animEffect transition="in" filter="box(in)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'fade_anim':
        extra = f'<p:animEffect transition="in" filter="checkerboard(across)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'glow_anim':
        extra = f'<p:animEffect transition="in" filter="circle(in)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'wipe_anim':
        extra = f'<p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_x</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav></p:tavLst></p:anim><p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid+1}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="1+#ppt_h/2"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav></p:tavLst></p:anim>'
    elif effect_type == 'blend_anim':
        extra = f'<p:animEffect transition="in" filter="blend()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'blur_anim':
        extra = f'<p:animEffect transition="in" filter="blur(in)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'misty_anim':
        extra = f'<p:animEffect transition="in" filter="misty()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'swizzle_anim':
        extra = f'<p:animEffect transition="in" filter="swizzle()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'dissolve_anim':
        extra = f'<p:animEffect transition="in" filter="dissolve()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'explode_anim':
        extra = f'<p:animEffect transition="in" filter="explode()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'grow_anim':
        extra = f'<p:animEffect transition="in" filter="grow()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'ripple_anim':
        extra = f'<p:animEffect transition="in" filter="ripple()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'reveal_anim':
        extra = f'<p:animEffect transition="in" filter="reveal()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'roll_anim':
        extra = f'<p:animEffect transition="in" filter="roll()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'shrink_anim':
        extra = f'<p:animEffect transition="in" filter="shrink()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'circle_anim':
        extra = f'<p:animEffect transition="in" filter="circle()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'teeter_anim':
        extra = f'<p:animEffect transition="in" filter="teeter()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    elif effect_type == 'typeWriter_anim':
        extra = f'<p:animEffect transition="in" filter="typeWriter()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
    else:
        extra = ''
    
    return subtype, extra
def create_background(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    sl.slide_width = SW; sl.slide_height = SH
    rect = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    rect.fill.solid(); rect.fill.fore_color.rgb = BG; rect.line.fill.background()
    return sl

def add_textbox(sl, left, top, width, height, text, size=18, color=None, bold=False, align=PP_ALIGN.LEFT):
    if color is None: color = LT
    tb = sl.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.font.size = Pt(size)
    p.font.color.rgb = color; p.font.bold = bold; p.alignment = align
    return tb

def add_card(sl, left, top, width, height, title, subtitle=''):
    rect = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    rect.fill.solid(); rect.fill.fore_color.rgb = CD; rect.line.color.rgb = GD; rect.line.width = Pt(1.5)
    add_textbox(sl, left+0.2, top+0.15, width-0.4, 0.5, title, 16, GD, True)
    if subtitle: add_textbox(sl, left+0.2, top+0.6, width-0.4, height-0.8, subtitle, 12, LT)
    return rect

def add_line(sl, left, top, width):
    line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(3))
    line.fill.solid(); line.fill.fore_color.rgb = GD; line.line.fill.background()
    return line

def animate_pptx(input_path, output_path=None, transitions=None, slide_groups=None, effects=None, method='com', style=None):
    """为 PowerPoint 添加动画效果。

    Args:
        input_path: 输入 PPTX 文件路径
        output_path: 输出文件路径（可选）
        transitions: 转场效果列表，如 ['fade', 'zoom', 'fade']
        slide_groups: 幻灯片分组（XML 模式使用）
        effects: 入场动画效果列表
        method: 'com' 或 'xml'
        style: 风格模板，如 '简洁商务'、'动态活力' 等
    """
    # 处理风格模板
    if style and style in EFFECT_STYLES:
        if effects is None:
            effects = EFFECT_STYLES[style]
        if transitions is None:
            transitions = TRANSITION_STYLES.get(style, ['fade'] * 3)
    else:
        if effects is None:
            effects = DEFAULT_EFFECTS
        if transitions is None:
            transitions = DEFAULT_TRANSITIONS

    if not os.path.exists(input_path): raise FileNotFoundError(f'Input file not found: {input_path}')
    if output_path is None:
        base, ext = os.path.splitext(input_path)
        output_path = base + '_animated' + ext
    if method == 'com' and HAS_WIN32COM:
        return _animate_with_com(input_path, output_path, transitions, slide_groups, effects)
    else:
        return _animate_with_xml(input_path, output_path, transitions, slide_groups, effects)

def _patch_gencache():
    """Patch win32com gencache to avoid CLSIDToClassMap errors."""
    try:
        import win32com.client.gencache as gencache_mod
        orig_add = gencache_mod.AddModuleToCache
        def patched_add(*args, **kwargs):
            try:
                return orig_add(*args, **kwargs)
            except AttributeError as e:
                if 'CLSIDToClassMap' in str(e):
                    return None
                raise
        gencache_mod.AddModuleToCache = patched_add
    except Exception:
        pass

def _animate_with_com(input_path, output_path, transitions=None, slide_groups=None, effects=None):
    """COM 自动化方式添加动画（需要安装 PowerPoint）。
    注意：COM 只支持 fade (3849) 转场，其他转场通过 XML 补充
    """
    _patch_gencache()
    _os = __import__("os")
    _time = __import__("time")
    _os.system("taskkill /F /IM POWERPNT.EXE 2>nul")
    _time.sleep(1)

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = True
    _time.sleep(1)

    ppt = app.Presentations.Open(_os.path.abspath(input_path))
    num_slides = ppt.Slides.Count
    
    # 处理转场参数
    if transitions is None:
        transitions = ['fade'] * num_slides
    elif isinstance(transitions, list) and transitions and isinstance(transitions[0], str):
        # 将转场名称转换为 COM ID
        transitions = [COM_TRANSITION_IDS.get(t, 3849) for t in transitions]
    
    anim_pool = [v for v in COM_ENTRANCE_IDS.values() if v is not None]
    TRIGGER_ONCLICK = 2

    # 先用 COM 设置入场动画（所有幻灯片统一使用 fade 转场 3849）
    for i in range(1, num_slides + 1):
        slide = ppt.Slides(i)
        seq = slide.TimeLine.MainSequence
        while seq.Count > 0:
            try: seq(1).Delete()
            except: break
        _time.sleep(0.2)
        try:
            slide.SlideShowTransition.EntryEffect = transitions[i-1] if i-1 < len(transitions) else 3849
            slide.SlideShowTransition.AdvanceOnClick = True
            slide.SlideShowTransition.AdvanceOnTime = False
            slide.SlideShowTransition.Duration = 1.0  # 转场持续时间1秒
        except: pass
        
        content_shapes = []
        for sh in slide.Shapes:
            try:
                if sh.Type == 17 and (sh.Name.startswith("Text") or sh.Name.startswith("TextBox")):
                    content_shapes.append(sh)
            except: pass
        for idx, shape in enumerate(content_shapes):
            if idx >= 8: break
            anim = seq.AddEffect(shape, anim_pool[idx % len(anim_pool)], TRIGGER_ONCLICK, 0)
            anim.Timing.Duration = 1.0
        print(f"Slide {i}: {len(content_shapes)} shapes animated")

    # 保存中间文件
    mid_path = _os.path.join(_os.path.dirname(output_path), '_temp_com.pptx')
    ppt.SaveAs(mid_path)
    try: ppt.Close()
    except: pass
    try: app.Quit()
    except: pass

    # 直接重命名输出文件
    if _os.path.exists(output_path):
        _os.remove(output_path)
    _os.rename(mid_path, output_path)
    
    return output_path

def inject_animations(content, slide_idx, transitions, slide_groups=None, effect_specs=None, effects=None):
    """注入动画到幻灯片XML——纯字符串操作，不依赖COM模板。"""
    # 添加转场动画
    if transitions is not None:
        # 确保 transitions 是列表
        if isinstance(transitions, str):
            transitions = [transitions] * 100
        # 如果 transitions 列表不够长，补齐
        while len(transitions) < slide_idx:
            transitions.append('fade')
        trans_type = transitions[slide_idx - 1]
        # 保存原始字符串用于查找 XML
        trans_type_str = trans_type
        if isinstance(trans_type, str):
            trans_type = TRANS_MAP.get(trans_type, 3849)
        
        content = re.sub(r'<p:transition[^>]*>.*?</p:transition>', '', content, flags=re.DOTALL)
        # 使用原始字符串查找 XML
        trans_xml = TRANSITION_XML.get(trans_type_str, '<p:fade/>')
        # 添加 Duration 和 AdvanceOnClick 属性（与COM版一致）
        trans_full = f'<p:transition durCtrlCnt="hold" advClick="1" advTime="3000">{trans_xml}</p:transition>'
        content = content.replace('</p:cSld>', '</p:cSld>' + trans_full)
    
    # 添加 timing 结构
    if '<p:timing>' not in content:
        shape_pattern = r'cNvPr id="([^"]+)" name="([^"]+)"'
        shape_matches = re.findall(shape_pattern, content)
        valid_spids = [int(sid) for sid, name in shape_matches if name.startswith('Text') or name.startswith('TextBox')]
        
        if valid_spids:
            # 构建完整的timing结构（与COM版一致）
            pars = []
            bldps = []
            cid = 13  # 从13开始，避免与tmRoot/mainSeq的id=1,2冲突
            
            for i, spid in enumerate(valid_spids[:8]):
                # 使用用户自定义动画或默认循环
                if effects is not None and i < len(effects):
                    preset_id = effects[i] if isinstance(effects[i], int) else EFFECT_MAP.get(effects[i], 1)
                else:
                    preset_id = [1, 2, 4, 5, 6, 7][i % 6]  # appear, checkerboard, fly, fade, glow, wipe
                
                if preset_id == 1:  # appear
                    subtype = 0
                    extra = f'<p:animEffect transition="in" filter="appear()"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
                elif preset_id == 2:  # checkerboard
                    subtype = 4
                    extra = f'<p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_x</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav></p:tavLst></p:anim><p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid+1}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="1+#ppt_h/2"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav></p:tavLst></p:anim>'
                    cid += 2
                elif preset_id == 4:  # fly
                    subtype = 16
                    extra = f'<p:animEffect transition="in" filter="box(in)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
                    cid += 1
                elif preset_id == 5:  # fade
                    subtype = 10
                    extra = f'<p:animEffect transition="in" filter="checkerboard(across)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
                    cid += 1
                elif preset_id == 6:  # glow
                    subtype = 16
                    extra = f'<p:animEffect transition="in" filter="circle(in)"><p:cBhvr><p:cTn id="{cid}" dur="500"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl></p:cBhvr></p:animEffect>'
                    cid += 1
                elif preset_id == 7:  # wipe
                    subtype = 4
                    extra = f'<p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_x</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_x"/></p:val></p:tav></p:tavLst></p:anim><p:anim calcmode="lin" valueType="num"><p:cBhvr additive="base"><p:cTn id="{cid+1}" dur="500" fill="hold"/><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>ppt_y</p:attrName></p:attrNameLst></p:cBhvr><p:tavLst><p:tav tm="0"><p:val><p:strVal val="1+#ppt_h/2"/></p:val></p:tav><p:tav tm="100000"><p:val><p:strVal val="#ppt_y"/></p:val></p:tav></p:tavLst></p:anim>'
                    cid += 2
                else:
                    subtype = 0
                    extra = ''
                
                # 使用 cid 作为 base_ctn_id，避免重复
                base_ctn_id = cid
                par = f'<p:par><p:cTn id="{base_ctn_id}" fill="hold"><p:stCondLst><p:cond delay="indefinite"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{base_ctn_id+1}" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:par><p:cTn id="{base_ctn_id+2}" presetID="{preset_id}" presetClass="entr" presetSubtype="{subtype}" fill="hold" grpId="0" nodeType="clickEffect"><p:stCondLst><p:cond delay="0"/></p:stCondLst><p:childTnLst><p:set><p:cBhvr><p:cTn id="{base_ctn_id+3}" dur="500" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn><p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="0" end="0"/></p:txEl></p:spTgt></p:tgtEl><p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst></p:cBhvr><p:to><p:strVal val="visible"/></p:to></p:set>{extra}</p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par></p:childTnLst></p:cTn></p:par>'
                pars.append(par)
                bldps.append(f'<p:bldP spid="{spid}" grpId="0" build="p"/>')
                cid += 4  # 每个 par 使用 4 个 cTn id
            
            # 构建完整的timing结构（与COM版一致）
            # COM 版结构: tmRoot -> childTnLst -> seq -> childTnLst -> mainSeq -> childTnLst -> [pars] -> /childTnLst -> /mainSeq -> /seq -> /childTnLst -> tmRoot
            # 注意：只有一个 </p:seq>，在 prevCondLst/nextCondLst 之后
            timing = f'<p:timing><p:tnLst><p:par><p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot"><p:childTnLst><p:seq concurrent="1" nextAc="seek"><p:cTn id="2" dur="indefinite" nodeType="mainSeq"><p:childTnLst>' + ''.join(pars) + '</p:childTnLst></p:cTn><p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst><p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst></p:seq></p:childTnLst></p:cTn></p:par></p:tnLst><p:bldLst>' + ''.join(bldps) + '</p:bldLst></p:timing>'
            
            # 找到 transition 标签的位置，将 timing 插入到 transition 后面
            trans_end = content.find('</p:transition>')
            if trans_end > 0:
                insert_pos = trans_end + len('</p:transition>')
                content = content[:insert_pos] + timing + content[insert_pos:]
            else:
                # 如果没有 transition，插入到 cSld 后面
                content = content.replace('</p:cSld>', '</p:cSld>' + timing)
    
    return content


def _animate_with_xml(input_path, output_path, transitions=None, slide_groups=None, effects=None, effect_specs=None):
    """XML 方式添加动画——纯字符串操作，不依赖 COM 模板。"""
    tmp_dir = tempfile.mkdtemp(prefix='ppt_animate_')
    try:
        # 解压输入文件
        with zipfile.ZipFile(input_path, 'r') as z:
            z.extractall(tmp_dir)
        
        slide_files = sorted([f for f in os.listdir(os.path.join(tmp_dir, 'ppt', 'slides')) if f.endswith('.xml')])

        num_slides = len(slide_files)
        if transitions is None:
            transitions = ['fade'] * num_slides
        for slide_file in slide_files:
            slide_idx = int(re.search(r'slide(\d+)', slide_file).group(1))
            slide_path = os.path.join(tmp_dir, 'ppt', 'slides', slide_file)
            with open(slide_path, 'r', encoding='utf-8') as f:
                content = f.read()
            content = inject_animations(content, slide_idx, transitions, slide_groups, effect_specs, effects)
            with open(slide_path, 'w', encoding='utf-8') as f:
                f.write(content)
        with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
            for root, dirs, files in os.walk(tmp_dir):
                for file in files:
                    file_path = os.path.join(root, file)
                    arcname = os.path.relpath(file_path, tmp_dir)
                    zip_out.write(file_path, arcname)
        
        print(f'Animated PPT saved to: {output_path} (XML method)')
        return output_path
    finally:
        shutil.rmtree(tmp_dir, ignore_errors=True)


def animate_existing_ppt(input_path, output_path=None, animation_type='click'):
    return animate_pptx(input_path, output_path, method='com')

def get_animation_preview(prs, slide_idx=1):
    if slide_idx < 1 or slide_idx > len(prs.slides):
        return None
    slide = prs.slides[slide_idx - 1]
    shapes = list(slide.shapes)
    groups = [[i+2] for i in range(len(shapes))]
    return {'slide_index': slide_idx, 'shape_count': len(shapes), 'groups': groups,
            'transitions': ['fade'] * len(groups)}

def get_animations_from_ppt(ppt_path):
    """从 PPT 中提取动画配置。"""
    anim_configs = []
    with zipfile.ZipFile(ppt_path, 'r') as z:
        for i in range(1, 100):
            slide_file = f'ppt/slides/slide{i}.xml'
            if slide_file not in z.namelist():
                break
            slide_content = z.read(slide_file).decode('utf-8')
            trans_match = re.search(r'entryEffect=\"(\d+)\"', slide_content)
            trans = int(trans_match.group(1)) if trans_match else None
            pars = re.findall(r'<[^>]+:par[^>]*>.*?</[^>]+:par>', slide_content, re.DOTALL)
            shape_effects = []
            for par in pars:
                spid_m = re.search(r'spid=\"(\d+)\"', par)
                filter_m = re.search(r'filter=\"([^\"]+)\"', par)
                preset_m = re.search(r'presetID=\"(\d+)\"', par)
                if spid_m:
                    spid = spid_m.group(1)
                    if filter_m:
                        effect = filter_m.group(1)
                        preset_id = preset_m.group(1) if preset_m else None
                    elif preset_m:
                        effect = f'preset:{preset_m.group(1)}'
                        preset_id = preset_m.group(1)
                    else:
                        effect = 'appear'
                        preset_id = None
                    shape_effects.append({'spid': spid, 'effect': effect, 'preset_id': preset_id})
            anim_configs.append({'slide': i, 'transition': trans, 'shape_effects': shape_effects})
    return anim_configs

def replicate_animations(source_ppt, target_ppt, output_path):
    """将源 PPT 的动画精确复刻到目标 PPT。"""
    source_configs = get_animations_from_ppt(source_ppt)
    effect_specs = []
    for config in source_configs:
        specs = [(se['spid'], se['effect'], se.get('preset_id')) for se in config['shape_effects']]
        effect_specs.append(specs)
    if os.path.abspath(target_ppt) == os.path.abspath(output_path):
        tmp = tempfile.mktemp(suffix='.pptx')
        shutil.copy2(target_ppt, tmp)
        try:
            return _animate_with_xml(tmp, output_path, effect_specs=effect_specs)
        finally:
            if os.path.exists(tmp):
                os.remove(tmp)
    else:
        return _animate_with_xml(target_ppt, output_path, effect_specs=effect_specs)

if __name__ == '__main__':
    import sys
    if len(sys.argv) > 1:
        input_file = sys.argv[1]
        output_file = input_file.replace('.pptx', '_animated.pptx')
        print(f'Processing {input_file}...')
        result = animate_pptx(input_file, output_file, method='com')
        print(f'Result: {result}')














